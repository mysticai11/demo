import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import math

BG       = RGBColor(0x05, 0x08, 0x10)
PANEL    = RGBColor(0x0C, 0x11, 0x1E)
BORDER   = RGBColor(0x1C, 0x29, 0x40)
TEXT_COL = RGBColor(0xEE, 0xF2, 0xFF)
MUTED    = RGBColor(0x4A, 0x63, 0x80)
TEAL     = RGBColor(0x00, 0xC4, 0x7D)
AMBER    = RGBColor(0xF5, 0xA6, 0x23)
BLUE     = RGBColor(0x3D, 0x8E, 0xF8)
RED      = RGBColor(0xE8, 0x39, 0x4A)
GOLD     = RGBColor(0xC8, 0xA8, 0x4B)

W_IN, H_IN = 13.33, 7.5
W = Inches(W_IN)
H = Inches(H_IN)
MARGIN = Inches(W_IN * 0.08)

def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width_pt=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width_pt:
            shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape

def set_opacity(shape, alpha_percent):
    solidFill = shape._element.find('.//' + qn('a:solidFill'))
    if solidFill is not None:
        srgb = solidFill.find(qn('a:srgbClr'))
        if srgb is None:
            srgb = solidFill.find(qn('a:schemeClr'))
        if srgb is not None:
            alpha_el = srgb.find(qn('a:alpha'))
            if alpha_el is None:
                alpha_el = etree.SubElement(srgb, qn('a:alpha'))
            val = int((alpha_percent / 100.0) * 100000)
            alpha_el.set('val', str(val))

def add_blob(slide, cx_in, cy_in, r_in, color, opacity=20):
    r = Inches(r_in)
    shape = slide.shapes.add_shape(9, Inches(cx_in)-r, Inches(cy_in)-r, r*2, r*2)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    set_opacity(shape, opacity)

def add_textbox(slide, text, left, top, width, height, font_name="Inter", font_size=24, bold=False, color=TEXT_COL, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_mappin(slide, cx_in, cy_in, size_in=0.45, color=TEAL):
    r = Inches(size_in * 0.4)
    cx, cy = Inches(cx_in), Inches(cy_in)
    circ = slide.shapes.add_shape(9, cx - r, cy - r*1.2, r*2, r*2)
    circ.fill.solid(); circ.fill.fore_color.rgb = color; circ.line.fill.background()
    tri = slide.shapes.add_shape(5, cx - r*0.8, cy - r*0.2, r*1.6, r*1.5)
    spPr = tri._element.find(qn('p:spPr'))
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is None: xfrm = etree.SubElement(spPr, qn('a:xfrm'))
    xfrm.set('rot', str(int(math.pi * 10800000)))
    tri.fill.solid(); tri.fill.fore_color.rgb = color; tri.line.fill.background()

def add_chartbar_icon(slide, cx_in, cy_in, size_in=0.45, color=TEAL):
    cx, cy = Inches(cx_in), Inches(cy_in)
    w = Inches(size_in * 0.2)
    h_max = Inches(size_in * 0.8)
    for i, h_mult in enumerate([0.4, 0.7, 1.0]):
        add_rect(slide, cx - Inches(size_in*0.4) + i*w*1.5, cy + (h_max - h_max*h_mult)/2, w, h_max*h_mult, fill_color=color)

def add_shield_icon(slide, cx_in, cy_in, size_in=0.45, color=TEAL):
    cx, cy = Inches(cx_in), Inches(cy_in)
    w = Inches(size_in * 0.8)
    h = Inches(size_in * 0.5)
    add_rect(slide, cx - w/2, cy - Inches(size_in*0.4), w, h, fill_color=color)
    tri = slide.shapes.add_shape(5, cx - w/2, cy + h - Inches(size_in*0.4), w, h)
    spPr = tri._element.find(qn('p:spPr'))
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is None: xfrm = etree.SubElement(spPr, qn('a:xfrm'))
    xfrm.set('rot', str(int(math.pi * 10800000)))
    tri.fill.solid(); tri.fill.fore_color.rgb = color; tri.line.fill.background()

def add_droplet_icon(slide, cx_in, cy_in, size_in=0.45, color=TEAL):
    r = Inches(size_in * 0.35)
    cx, cy = Inches(cx_in), Inches(cy_in)
    circ = slide.shapes.add_shape(9, cx - r, cy, r*2, r*2)
    circ.fill.solid(); circ.fill.fore_color.rgb = color; circ.line.fill.background()
    tri = slide.shapes.add_shape(5, cx - r, cy - r*1.2, r*2, r*1.8)
    tri.fill.solid(); tri.fill.fore_color.rgb = color; tri.line.fill.background()

def add_clock_icon(slide, cx_in, cy_in, size_in=0.45, color=TEAL):
    r = Inches(size_in * 0.4)
    cx, cy = Inches(cx_in), Inches(cy_in)
    circ = slide.shapes.add_shape(9, cx - r, cy - r, r*2, r*2)
    circ.fill.background()
    circ.line.color.rgb = color
    circ.line.width = Pt(3)
    add_rect(slide, cx - Pt(1.5), cy - r*0.6, Pt(3), r*0.6, fill_color=color)
    add_rect(slide, cx, cy - Pt(1.5), r*0.5, Pt(3), fill_color=color)

def add_icon_badge(slide, cx_in, cy_in, color, icon_type):
    bg = slide.shapes.add_shape(9, Inches(cx_in - 0.4), Inches(cy_in - 0.4), Inches(0.8), Inches(0.8))
    bg.fill.solid(); bg.fill.fore_color.rgb = color; bg.line.fill.background()
    set_opacity(bg, 15)
    if icon_type == 'pin': add_mappin(slide, cx_in, cy_in, 0.45, color)
    elif icon_type == 'chart': add_chartbar_icon(slide, cx_in, cy_in, 0.45, color)
    elif icon_type == 'shield': add_shield_icon(slide, cx_in, cy_in, 0.45, color)
    elif icon_type == 'droplet': add_droplet_icon(slide, cx_in, cy_in, 0.45, color)
    elif icon_type == 'clock': add_clock_icon(slide, cx_in, cy_in, 0.45, color)
    else:
        add_rect(slide, Inches(cx_in)-Inches(0.1), Inches(cy_in)+Inches(0.1), Inches(0.2), Inches(0.05), fill_color=color)

def build_presentation():
    prs = new_prs()
    layout = prs.slide_layouts[6]
    
    # 1. Title
    s = prs.slides.add_slide(layout); set_bg(s, BG)
    add_blob(s, 3, 2, 2.5, TEAL, 10)
    add_blob(s, 10, 5, 3.0, BLUE, 10)
    add_blob(s, 10, 2, 2.0, AMBER, 10)
    add_blob(s, 3, 5, 2.5, RED, 10)
    add_textbox(s, "Predictive Risk Intelligence for Metabolic Screening in Diabetes", MARGIN, Inches(2), W - 2*MARGIN, Inches(1.5), font_name="Poppins", font_size=44, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, "LMSIS — Latent Metabolic State Inference System", MARGIN, Inches(3.5), W - 2*MARGIN, Inches(1), font_size=24, color=TEAL, align=PP_ALIGN.CENTER)
    presenters = "Karamjeet Singh (5244) • Sarandeep Singh (5244) • Ubaid Ahmad (5255) • Aiman Younus Dar (222176)"
    add_textbox(s, presenters, MARGIN, Inches(5), W - 2*MARGIN, Inches(0.5), font_size=16, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, "SHMM Government Degree College Anantnag", MARGIN, Inches(5.4), W - 2*MARGIN, Inches(0.5), font_size=18, color=TEAL, align=PP_ALIGN.CENTER)
    add_textbox(s, "June 2026", MARGIN, Inches(5.8), W - 2*MARGIN, Inches(0.5), font_size=16, color=MUTED, align=PP_ALIGN.CENTER)

    # 2. The Paradox
    s = prs.slides.add_slide(layout); set_bg(s, BG)
    add_textbox(s, "Normal BMI. Normal cholesterol. No symptoms.", MARGIN, Inches(2.5), W - 2*MARGIN, Inches(1), font_name="Poppins", font_size=36, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, "Their doctor sent them home.", MARGIN, Inches(3.5), W - 2*MARGIN, Inches(1), font_name="Poppins", font_size=36, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, "They had severe liver disease.", MARGIN, Inches(4.5), W - 2*MARGIN, Inches(1), font_name="Poppins", font_size=36, bold=True, color=RED, align=PP_ALIGN.CENTER)

    # 3. Why BMI Misses This
    s = prs.slides.add_slide(layout); set_bg(s, BG)
    add_mappin(s, W_IN - 1.0, 6.8, 0.4, AMBER)
    fw, fh = Inches(1.5), Inches(3.0)
    for i in range(2):
        cx = W/2 - Inches(2.5) if i==0 else W/2 + Inches(1.0)
        add_rect(s, cx, Inches(1.5), fw, fh, line_color=BORDER, line_width_pt=3)
        add_rect(s, cx + fw/2 - Inches(0.3), Inches(1.2), Inches(0.6), Inches(0.6), line_color=BORDER, line_width_pt=3)
        add_textbox(s, "BMI 22", cx, Inches(4.8), fw, Inches(0.5), font_name="JetBrains Mono", font_size=24, align=PP_ALIGN.CENTER)
        if i == 1:
            add_blob(s, cx/Inches(1) + 0.75, 3.0, 0.4, AMBER, 80)
            add_rect(s, cx + Inches(1.2), Inches(3.0), Inches(0.5), Pt(2), fill_color=AMBER)
    add_textbox(s, "BMI measures weight. It can't see where fat is hiding.", MARGIN, H - Inches(1.5), W - 2*MARGIN, Inches(0.5), font_size=24, align=PP_ALIGN.CENTER)

    # 4. The Problem
    s = prs.slides.add_slide(layout); set_bg(s, BG)
    texts = [
        ("Standard scores lose accuracy at normal weight", 'chart', AMBER),
        ("One major score gets it backwards entirely", 'droplet', RED),
        ("AI models that work, but explain nothing", 'shield', BLUE),
        ("Confidence estimates that fail the riskiest patients", 'clock', RED)
    ]
    for i, (txt, icon, col) in enumerate(texts):
        row, col_idx = i // 2, i % 2
        left = MARGIN + col_idx * (W/2 - MARGIN + Inches(0.2))
        top = Inches(1.5) + row * Inches(2.5)
        add_rect(s, left, top, W/2 - MARGIN - Inches(0.4), Inches(2), fill_color=PANEL, line_color=BORDER, line_width_pt=1)
        add_icon_badge(s, left/Inches(1) + 0.8, top/Inches(1) + 1.0, col, icon)
        add_textbox(s, txt, left + Inches(1.6), top + Inches(0.5), W/2 - MARGIN - Inches(2.2), Inches(1), font_size=24)

    # 5. The Gap
    s = prs.slides.add_slide(layout); set_bg(s, BG)
    items = [
        "Built specifically for normal-weight patients",
        "Risk scores tied to real medical tests",
        "Reliable confidence — even for the riskiest patients",
        "Plain explanations, not a black box"
    ]
    for i, txt in enumerate(items):
        y = Inches(1.5) + i * Inches(1.0)
        add_icon_badge(s, MARGIN/Inches(1) + 0.5, y/Inches(1) + 0.3, TEAL, 'shield')
        add_textbox(s, txt, MARGIN + Inches(1.2), y, W - 2*MARGIN, Inches(0.6), font_size=28)
    add_textbox(s, "No existing system does all four. This one does.", MARGIN, Inches(6.0), W - 2*MARGIN, Inches(0.5), font_size=22, color=MUTED, align=PP_ALIGN.CENTER)

    # 6. How It Works
    s = prs.slides.add_slide(layout); set_bg(s, BG)
    steps = [
        ("Blood test values", 'droplet', TEAL),
        ("Pattern read as a whole", 'chart', BLUE),
        ("Two risk scores", 'chart', AMBER),
        ("Patient placed on a map", 'pin', RED),
        ("Confidence + route to safety", 'shield', TEAL)
    ]
    w_step = Inches(2.0)
    for i, (txt, icon, col) in enumerate(steps):
        x = MARGIN + i * (w_step + Inches(0.3))
        add_rect(s, x, Inches(3.0), w_step, Inches(1.8), fill_color=PANEL, line_color=BORDER, line_width_pt=1)
        add_icon_badge(s, x/Inches(1) + 1.0, 3.4, col, icon)
        add_textbox(s, txt, x, Inches(3.9), w_step, Inches(0.8), font_size=18, align=PP_ALIGN.CENTER)
        if i < 4:
            add_rect(s, x + w_step + Inches(0.05), Inches(3.8), Inches(0.2), Pt(2), fill_color=BORDER)
    add_textbox(s, "Implemented as an identifiable generative model, anchored to real clinical measurements — full detail in the written report.", MARGIN, H - Inches(1.0), W - 2*MARGIN, Inches(0.5), font_size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # 7. Live Demonstration
    s = prs.slides.add_slide(layout); set_bg(s, BG)
    add_textbox(s, "Live demonstration", MARGIN, Inches(3.0), W - 2*MARGIN, Inches(1.5), font_name="Poppins", font_size=48, bold=True, align=PP_ALIGN.CENTER)

    # 8. Why Existing Tools Fail
    s = prs.slides.add_slide(layout); set_bg(s, BG)
    add_textbox(s, "One widely used clinical score ranks sicker patients as healthier.", MARGIN, H - Inches(1.5), W - 2*MARGIN, Inches(0.5), font_size=24, align=PP_ALIGN.CENTER)
    bars = [("LMSIS", 0.8, TEAL), ("Score A", 0.6, MUTED), ("Score B", 0.5, MUTED), ("Score C", 0.4, MUTED), ("Flawed Score", -0.4, RED)]
    for i, (name, val, col) in enumerate(bars):
        y = Inches(1.5) + i * Inches(0.8)
        add_textbox(s, name, MARGIN, y, Inches(2), Inches(0.5), font_size=20, align=PP_ALIGN.RIGHT)
        bw = Inches(abs(val) * 4)
        if val > 0:
            add_rect(s, MARGIN + Inches(2.2), y + Inches(0.1), bw, Inches(0.4), fill_color=col)
            add_textbox(s, f"{val:.2f}", MARGIN + Inches(2.3) + bw, y, Inches(1), Inches(0.5), font_name="JetBrains Mono", font_size=20, color=col)
        else:
            add_rect(s, MARGIN + Inches(2.2) - bw, y + Inches(0.1), bw, Inches(0.4), fill_color=col)
            add_textbox(s, f"{val:.2f}", MARGIN + Inches(1.0) - bw, y, Inches(1), Inches(0.5), font_name="JetBrains Mono", font_size=20, color=col)
    add_rect(s, MARGIN + Inches(2.2), Inches(1.3), Pt(2), Inches(4.0), fill_color=BORDER)

    # 9. Trustworthy Confidence Ranges
    s = prs.slides.add_slide(layout); set_bg(s, BG)
    add_textbox(s, "Standard methods miss the riskiest patients. This fixes that.", MARGIN, H - Inches(1.5), W - 2*MARGIN, Inches(0.5), font_size=24, align=PP_ALIGN.CENTER)
    h_before, h_after = Inches(2.0), Inches(3.5)
    y_base = Inches(5.0)
    add_rect(s, W/2 - Inches(2), y_base - h_before, Inches(1.5), h_before, fill_color=MUTED)
    add_textbox(s, "81.6%", W/2 - Inches(2), y_base - h_before - Inches(0.6), Inches(1.5), Inches(0.5), font_name="JetBrains Mono", font_size=32, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, "Before", W/2 - Inches(2), y_base + Inches(0.1), Inches(1.5), Inches(0.5), font_size=24, color=MUTED, align=PP_ALIGN.CENTER)
    add_rect(s, W/2 + Inches(0.5), y_base - h_after, Inches(1.5), h_after, fill_color=TEAL)
    add_textbox(s, "90.4%", W/2 + Inches(0.5), y_base - h_after - Inches(0.6), Inches(1.5), Inches(0.5), font_name="JetBrains Mono", font_size=32, color=TEAL, align=PP_ALIGN.CENTER, bold=True)
    add_textbox(s, "After", W/2 + Inches(0.5), y_base + Inches(0.1), Inches(1.5), Inches(0.5), font_size=24, color=TEXT_COL, align=PP_ALIGN.CENTER)

    # 10. Tested Beyond Training Data
    s = prs.slides.add_slide(layout); set_bg(s, BG)
    add_blob(s, 6.5, 4, 3, TEAL, 5)
    panels = [
        ("Tested two years later, on patients never seen before", "ρ = 0.583 ± 0.027", 'clock', TEAL),
        ("Tested on a separate population group", "ρ = 0.582 (N=355)", 'shield', BLUE)
    ]
    for i, (txt, res, icon, col) in enumerate(panels):
        x = MARGIN + i * (W/2 - MARGIN + Inches(0.2))
        add_rect(s, x, Inches(2.0), W/2 - MARGIN - Inches(0.4), Inches(3.0), fill_color=PANEL, line_color=BORDER, line_width_pt=1)
        add_icon_badge(s, x/Inches(1) + 0.8, 2.6, col, icon)
        add_textbox(s, txt, x + Inches(0.4), Inches(3.2), W/2 - MARGIN - Inches(1.2), Inches(1), font_size=20, color=MUTED)
        add_textbox(s, res, x + Inches(0.4), Inches(4.0), W/2 - MARGIN - Inches(1.2), Inches(1.0), font_name="JetBrains Mono", font_size=32, bold=True, color=col)

    # 11. What We Built
    s = prs.slides.add_slide(layout); set_bg(s, BG)
    lines = [
        ("Risk scores mathematically anchored to real clinical measurements, not left as an unexplainable black box", 'pin', TEAL),
        ("Reliable confidence ranges specifically for the highest-risk patients, where standard methods are shown to fail", 'shield', BLUE),
        ("Plain mathematical formulas explaining how the model behaves", 'chart', AMBER),
        ("A complete, tested, working system — not just a research idea", 'shield', TEAL)
    ]
    for i, (txt, icon, col) in enumerate(lines):
        y = Inches(1.5) + i * Inches(1.2)
        add_rect(s, MARGIN, y, W - 2*MARGIN, Inches(1.0), fill_color=PANEL, line_color=BORDER, line_width_pt=1)
        add_icon_badge(s, MARGIN/Inches(1) + 0.6, y/Inches(1) + 0.5, col, icon)
        add_textbox(s, txt, MARGIN + Inches(1.3), y + Inches(0.2), W - 2*MARGIN - Inches(1.6), Inches(0.8), font_size=18)

    # 12. Questions You Might Be Asking
    s = prs.slides.add_slide(layout); set_bg(s, BG)
    qa = [
        ("Is this ready for real patients?", "No — a research prototype, not a diagnostic device."),
        ("Why not just use a simpler model?", "Because it would give no explainable structure or per-group confidence."),
        ("Could this just be overfit to one dataset?", "No — tested on data collected two years later that it never saw."),
        ("Why such a wide range on the national estimate?", "It honestly reflects a small subgroup sample size, not a flawed model.")
    ]
    for i, (q, a) in enumerate(qa):
        row, col_idx = i // 2, i % 2
        x = MARGIN + col_idx * (W/2 - MARGIN + Inches(0.2))
        y = Inches(1.5) + row * Inches(2.5)
        add_rect(s, x, y, W/2 - MARGIN - Inches(0.4), Inches(2.2), fill_color=PANEL, line_color=BORDER, line_width_pt=1)
        add_textbox(s, q, x + Inches(0.3), y + Inches(0.3), W/2 - MARGIN - Inches(1.0), Inches(0.8), font_name="Poppins", font_size=20, bold=True)
        add_textbox(s, a, x + Inches(0.3), y + Inches(1.2), W/2 - MARGIN - Inches(1.0), Inches(0.8), font_size=18, color=MUTED)

    # 13. Honest Limits and What's Next
    s = prs.slides.add_slide(layout); set_bg(s, BG)
    add_rect(s, MARGIN, Inches(3.75), W - 2*MARGIN, Pt(2), fill_color=BORDER)
    add_textbox(s, "What we don't yet know", MARGIN, Inches(1.0), W - 2*MARGIN, Inches(0.5), font_size=18, color=MUTED)
    top_items = ["Built from a single national data source", "A one-time snapshot, not tracked over time", "Not a diagnostic device"]
    for i, t in enumerate(top_items):
        add_textbox(s, "• " + t, MARGIN + Inches(0.5), Inches(1.8) + i * Inches(0.6), W - 2*MARGIN, Inches(0.5), font_size=22)
    add_textbox(s, "What's next", MARGIN, Inches(4.2), W - 2*MARGIN, Inches(0.5), font_size=18, color=TEAL)
    bot_items = ["Test on completely independent data", "Track real patients over time", "Expand to more population groups"]
    for i, t in enumerate(bot_items):
        add_textbox(s, "• " + t, MARGIN + Inches(0.5), Inches(5.0) + i * Inches(0.6), W - 2*MARGIN, Inches(0.5), font_size=22)

    # 14. Closing
    s = prs.slides.add_slide(layout); set_bg(s, BG)
    add_textbox(s, "Normal BMI is not metabolic health.", MARGIN, Inches(2.5), W - 2*MARGIN, Inches(0.8), font_name="Poppins", font_size=38, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, "We can show you where the patient actually is — and the route back.", MARGIN, Inches(3.3), W - 2*MARGIN, Inches(0.8), font_name="Poppins", font_size=38, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    
    # badges
    add_rect(s, W/2 - Inches(2.5), Inches(4.5), Inches(1.3), Inches(0.4), fill_color=PANEL, line_color=BORDER, line_width_pt=1)
    add_textbox(s, "Tested", W/2 - Inches(2.5), Inches(4.55), Inches(1.3), Inches(0.4), font_size=16, color=MUTED, align=PP_ALIGN.CENTER)
    add_rect(s, W/2 - Inches(0.8), Inches(4.5), Inches(1.6), Inches(0.4), fill_color=PANEL, line_color=BORDER, line_width_pt=1)
    add_textbox(s, "Reproducible", W/2 - Inches(0.8), Inches(4.55), Inches(1.6), Inches(0.4), font_size=16, color=MUTED, align=PP_ALIGN.CENTER)
    add_rect(s, W/2 + Inches(1.2), Inches(4.5), Inches(1.3), Inches(0.4), fill_color=PANEL, line_color=BORDER, line_width_pt=1)
    add_textbox(s, "Open Source", W/2 + Inches(1.2), Inches(4.55), Inches(1.3), Inches(0.4), font_size=16, color=MUTED, align=PP_ALIGN.CENTER)
    
    add_textbox(s, "Thank you. Questions?", MARGIN, Inches(6.0), W - 2*MARGIN, Inches(0.8), font_size=22, color=MUTED, align=PP_ALIGN.CENTER)
    
    prs.save("LMSIS_Demo_Presentation.pptx")

if __name__ == "__main__":
    build_presentation()
